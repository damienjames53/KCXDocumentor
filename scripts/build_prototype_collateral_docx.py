#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from copy import deepcopy

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor

WORKSPACE = Path(__file__).resolve().parents[1]
OUT_DIR = WORKSPACE / "docs" / "prototype-collateral"
BRAND_IMAGES = WORKSPACE / "assets" / "branding" / "images"
BRAND_TEMPLATES = WORKSPACE / "assets" / "branding" / "templates"
PPT_TEMPLATE = BRAND_TEMPLATES / "Keycentrix Powerpoint Template - Internal and Support 1.pptx"
SAMPLE_GUIDE = WORKSPACE / "artifacts" / "generated" / "blink-rx-training-part-1-112125-4eba47fa6aa0" / "user_guide.anthropic.docx"

PRIMARY = RGBColor(0x1C, 0x75, 0xBC)
GREEN = RGBColor(0x8C, 0xC6, 0x3F)
DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BLUE = "EAF3FB"
LIGHT_GRAY = "F3F4F6"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


@dataclass(frozen=True)
class DocSpec:
    title: str
    subtitle: str
    document_id: str
    status: str
    owner: str
    audience: str
    purpose: str
    scope: str
    sections: list[tuple[str, list[str] | list[list[str]]]]
    output_name: str


def _first_logo() -> Path | None:
    for name in ("keycentrix-template-logo.png", "keycentrix-full-logo.png", "sticky-logo.png"):
        path = BRAND_IMAGES / name
        if path.exists():
            return path
    return None


def _set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def _set_cell_margins(cell, top=80, start=120, bottom=80, end=120) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for m, v in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(v))
        node.set(qn("w:type"), "dxa")


def _set_cell_border(cell, color="D7DDE0", size="4") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = tc_pr.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        tc_pr.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        node = borders.find(qn(f"w:{edge}"))
        if node is None:
            node = OxmlElement(f"w:{edge}")
            borders.append(node)
        node.set(qn("w:val"), "single")
        node.set(qn("w:sz"), size)
        node.set(qn("w:color"), color)


def _style_document(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = DARK
    normal.paragraph_format.space_after = Pt(6)

    for style_name, size, color in (("Heading 1", 16, PRIMARY), ("Heading 2", 12.5, DARK), ("Heading 3", 11, PRIMARY)):
        style = doc.styles[style_name]
        style.font.name = "Calibri"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = color
        style.paragraph_format.space_before = Pt(12)
        style.paragraph_format.space_after = Pt(5)


def _header_footer(doc: Document, title: str) -> None:
    section = doc.sections[0]
    header = section.header.paragraphs[0]
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = header.add_run(title)
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.color.rgb = MUTED

    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = footer.add_run("keycentrix | KCXDocumentor")
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.color.rgb = MUTED


def _table(doc: Document, headers: list[str], rows: list[list[str]], widths: list[float] | None = None) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for idx, header in enumerate(headers):
        cell = table.rows[0].cells[idx]
        cell.text = header
        _set_cell_shading(cell, "1C75BC")
        _set_cell_margins(cell)
        _set_cell_border(cell)
        for run in cell.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = WHITE
    for row_values in rows:
        row = table.add_row()
        for idx, value in enumerate(row_values):
            cell = row.cells[idx]
            cell.text = value
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _set_cell_shading(cell, LIGHT_GRAY if idx == 0 else "FFFFFF")
            _set_cell_margins(cell)
            _set_cell_border(cell)
            for run in cell.paragraphs[0].runs:
                run.font.name = "Calibri"
                run.font.size = Pt(9.5)
                run.font.color.rgb = DARK
                if idx == 0:
                    run.font.bold = True
    if widths:
        for row in table.rows:
            for idx, width in enumerate(widths):
                row.cells[idx].width = Inches(width)
    doc.add_paragraph()


def _callout(doc: Document, title: str, body: str, fill: str = LIGHT_BLUE) -> None:
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    cell = table.rows[0].cells[0]
    _set_cell_shading(cell, fill)
    _set_cell_margins(cell, top=120, start=160, bottom=120, end=160)
    _set_cell_border(cell, color="8CC63F", size="6")
    p = cell.paragraphs[0]
    r = p.add_run(title)
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    p = cell.add_paragraph()
    r = p.add_run(body)
    r.font.name = "Calibri"
    r.font.size = Pt(10)
    r.font.color.rgb = DARK
    doc.add_paragraph()


def _bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)


def _numbered(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Number")
        p.add_run(item)


def _cover(doc: Document, spec: DocSpec) -> None:
    logo = _first_logo()
    if logo:
        p = doc.add_paragraph()
        p.add_run().add_picture(str(logo), width=Inches(2.25))

    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(26)
    run = p.add_run(spec.title)
    run.font.name = "Calibri"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = PRIMARY

    p = doc.add_paragraph()
    run = p.add_run(spec.subtitle)
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED

    _callout(
        doc,
        "Recipe alignment",
        "Built as a keycentrix internal artifact using the DamienDev document recipe structure: control metadata, version history, purpose, scope, audience, governance, process detail, and appendices.",
        fill="EAF3FB",
    )

    _table(
        doc,
        ["Metadata Field", "Document Value"],
        [
            ["Document ID", spec.document_id],
            ["Version", "1.0"],
            ["Effective Date", date.today().isoformat()],
            ["Status", spec.status],
            ["Owner", spec.owner],
            ["Audience", spec.audience],
            ["Build Standard", "DamienDev document recipe"],
        ],
        [2.0, 4.8],
    )
    doc.add_page_break()


def _build_doc(spec: DocSpec) -> Path:
    doc = Document()
    doc.core_properties.title = spec.title
    doc.core_properties.author = "keycentrix"
    doc.core_properties.subject = spec.subtitle
    doc.core_properties.comments = "Built with the DamienDev document recipe pattern."
    _style_document(doc)
    _header_footer(doc, spec.title)
    _cover(doc, spec)

    doc.add_heading("Document Control", level=1)
    _table(
        doc,
        ["Control", "Value"],
        [
            ["Document Title", spec.title],
            ["Document ID", spec.document_id],
            ["Status", spec.status],
            ["Owner", spec.owner],
            ["Audience", spec.audience],
        ],
    )

    doc.add_heading("Version History", level=1)
    _table(doc, ["Version", "Date", "Summary"], [["1.0", date.today().isoformat(), "Initial collateral package aligned to the DamienDev document recipe."]])

    doc.add_heading("Purpose", level=1)
    doc.add_paragraph(spec.purpose)
    doc.add_heading("Executive Summary", level=1)
    _callout(
        doc,
        "What this document is for",
        "This file is part of the KCXDocumentor executive package and is intended to be readable as a standalone decision-support artifact.",
        fill="F3F4F6",
    )
    doc.add_heading("Scope", level=1)
    doc.add_paragraph(spec.scope)
    doc.add_heading("Intended Audience", level=1)
    doc.add_paragraph(spec.audience)
    doc.add_heading("Governing Principles", level=1)
    _bullets(
        doc,
        [
            "Keep raw recordings and local processing artifacts on the workstation.",
            "Send compact, reviewer-curated prompt context through the authenticated AI route.",
            "Keep human review in front of customer-facing publication.",
            "Track AI usage centrally by user, document, token count, page count, and estimated cost.",
        ],
    )

    for title, content in spec.sections:
        doc.add_heading(title, level=1)
        if content and isinstance(content[0], list):
            rows = content  # type: ignore[assignment]
            _table(doc, rows[0], rows[1:])  # type: ignore[arg-type]
        else:
            _bullets(doc, content)  # type: ignore[arg-type]

    output = OUT_DIR / spec.output_name
    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(output)
    return output


def _doc_specs() -> list[DocSpec]:
    return [
        DocSpec(
            title="KCXDocumentor - Executive Brief",
            subtitle="AI-assisted documentation from workflow recordings",
            document_id="KCXDOC-BRIEF-001",
            status="Internal executive review",
            owner="keycentrix Engineering and Product",
            audience="Executive Team, Product, Training, Implementation, and Engineering",
            purpose="Summarize the business value, operating model, and validation path for KCXDocumentor as an internal documentation acceleration capability.",
            scope="Covers the current application behavior, architecture, spend visibility, pilot validation approach, and principal risks for internal decision-making.",
            sections=[
                ("Business Problem", [
                    "Long workflow recordings require manual review before they become useful guides.",
                    "Business analysts and trainers spend time finding steps, rewriting narration, selecting screenshots, and formatting Word documents.",
                    "Document quality varies by author, recording quality, and available transcript detail.",
                ]),
                ("Current Capability", [
                    "Imports recordings and optional transcripts through the local web console.",
                    "Processes video, audio, transcript, OCR, and candidate frames locally.",
                    "Lets reviewers approve, reject, or capture screenshots before guide creation.",
                    "Creates a DOCX guide through an authenticated Azure Function proxy to Anthropic.",
                    "Persists AI Spend in Cosmos DB by user, document, token count, page count, and estimated cost.",
                ]),
                ("Architecture Summary", [
                    ["Layer", "Current Path", "Value"],
                    ["Local workstation", "Dockerized app with mapped folders", "Keeps recordings and artifacts under workstation control"],
                    ["Media processing", "FFmpeg, Whisper, OCR, frame review", "Creates compact procedure context from long recordings"],
                    ["Authentication", "Microsoft Entra with MSAL + PKCE", "Restricts use to approved users"],
                    ["AI route", "Azure Function proxy", "Keeps provider key server-side and validates user tokens"],
                    ["Reporting", "Cosmos DB AI Spend", "Preserves group-level usage visibility"],
                ]),
                ("Validation Measures", [
                    "Time to first usable draft compared with manual documentation.",
                    "Reviewer edits required before a guide can be shared.",
                    "Screenshot relevance and step alignment.",
                    "Absence of internal AI, prompt, or pipeline language in the reader-facing document.",
                    "Cost per guide and cost per generated page.",
                ]),
                ("Recommended Next Steps", [
                    "Run representative BA and trainer recordings through the application.",
                    "Compare transcript-backed runs against local Whisper-only runs.",
                    "Capture reviewer feedback on screenshot selection and guide clarity.",
                    "Define acceptance criteria for publishable documentation.",
                    "Confirm whether the next release should add native recording capture or continue with imported recordings.",
                ]),
            ],
            output_name="KCXDocumentor - Executive Brief.docx",
        ),
        DocSpec(
            title="KCXDocumentor - Demo Script and Storyboard",
            subtitle="Executive walkthrough script for the internal documentation workflow",
            document_id="KCXDOC-DEMO-001",
            status="Internal demonstration ready",
            owner="keycentrix Product and Training",
            audience="Executive sponsors, product leadership, trainers, business analysts, and implementation leads",
            purpose="Provide a concise 3 to 5 minute demo path that shows how KCXDocumentor turns a local recording into a reviewable Word guide.",
            scope="Covers the recommended demo objective, production notes, storyboard, narration points, and close-out message for the executive review.",
            sections=[
                ("Demo Objective", [
                    "Show the working flow from imported recording to DOCX output.",
                    "Emphasize local processing, authenticated AI access, reviewer control, and AI Spend visibility.",
                    "Keep the demo centered on the documentation workflow rather than implementation details.",
                ]),
                ("Production Notes", [
                    "Use a real source recording and show the import process.",
                    "State that raw video stays local while compact reviewed context is sent through the authenticated AI route.",
                    "Show frame approval and rejection because it is the primary quality gate.",
                    "Close with the validation ask for trainers and business analysts.",
                ]),
                ("Storyboard", [
                    ["Segment", "Visual", "Narration Focus"],
                    ["Opening", "KCXDocumentor workspace", "Recorded walkthroughs should not require hours of manual documentation work."],
                    ["Import", "Recording and transcript import controls", "A reviewer imports the source recording and optional transcript."],
                    ["Process", "Processing progress state", "Local media processing creates trace data and candidate screenshots."],
                    ["Review", "Frame reviewer", "The reviewer rejects overlays and captures better screenshots when needed."],
                    ["Create Guide", "AI creation progress", "Compact reviewed context is sent through the authenticated proxy."],
                    ["Download", "Download DOCX", "The result is a Word guide ready for human review."],
                    ["AI Spend", "AI Spend page", "Leadership can see documents, tokens, pages, costs, and user attribution."],
                    ["Close", "Validation summary", "The next step is representative internal testing with BAs and trainers."],
                ]),
                ("Close Message", [
                    "KCXDocumentor is ready for internal hands-on validation.",
                    "The review should focus on output quality, screenshot relevance, reviewer workflow, cost per finished page, and readiness for a broader rollout.",
                ]),
            ],
            output_name="KCXDocumentor - Demo Script and Storyboard.docx",
        ),
        DocSpec(
            title="KCXDocumentor - Windows Setup Guide",
            subtitle="Internal workstation setup guide for the Dockerized KCXDocumentor build",
            document_id="KCXDOC-SETUP-001",
            status="Internal tester setup ready",
            owner="keycentrix Engineering",
            audience="Business analysts, trainers, implementation specialists, documentation reviewers, and IT support users",
            purpose="Explain how an internal tester prepares a Windows workstation, starts KCXDocumentor, imports recordings, creates a guide, and reviews output.",
            scope="Covers Docker Desktop prerequisites, mapped folders, local environment settings, GHCR sign-in, startup, health checks, import flow, guide creation, AI Spend, and troubleshooting.",
            sections=[
                ("Prerequisites", [
                    "Windows 10 or Windows 11 workstation.",
                    "Docker Desktop installed, running, and set to Linux containers.",
                    "Internet access for first start unless Whisper has already been preseeded.",
                    "Microsoft account allowed to sign in to KCXDocumentor.",
                    "GitHub account or token with permission to pull the private container package.",
                ]),
                ("Folder Layout", [
                    ["Folder", "Purpose"],
                    ["C:\\KCXDocumentor\\samples\\raw", "Imported source videos and optional transcript files."],
                    ["C:\\KCXDocumentor\\samples\\processed", "Local processing sessions."],
                    ["C:\\KCXDocumentor\\artifacts", "Generated DOCX files and QA artifacts."],
                    ["C:\\KCXDocumentor\\external\\whisper", "Runtime Whisper binaries and models."],
                ]),
                ("Start The Application", [
                    "Copy .env.example to .env and confirm the mapped Windows folders.",
                    "Sign in to GitHub Container Registry with a token that has read:packages.",
                    "Run docker compose pull from the setup folder.",
                    "Run docker compose up -d.",
                    "Open http://127.0.0.1:8765 after the container reports healthy.",
                ]),
                ("Create A Guide", [
                    "Sign in with the approved Microsoft account.",
                    "Import the recording and optional transcript through the browser UI.",
                    "Choose Teams Recording when the video includes Teams overlays.",
                    "Process the recording and wait for the processing progress indicator to finish.",
                    "Review screenshots, reject irrelevant frames, and add better frames from the video when needed.",
                    "Choose Create Guide and wait for the AI creation progress indicator to complete.",
                    "Download the generated DOCX and review it before sharing.",
                ]),
                ("Troubleshooting", [
                    ["Issue", "What To Check"],
                    ["Docker cannot connect", "Confirm Docker Desktop is running and fully started."],
                    ["The app does not open", "Confirm docker compose ps shows port 8765."],
                    ["Sign-in fails", "Confirm the app registration allows http://127.0.0.1:8765/ and the user is assigned."],
                    ["No recordings appear", "Confirm the source video was imported or placed in the mapped raw folder."],
                    ["Whisper is unavailable", "Check container logs and confirm the Whisper folder is writable."],
                    ["Guide creation fails", "Review Latest Activity and confirm the Azure Function API is reachable."],
                ]),
            ],
            output_name="KCXDocumentor - Windows Setup Guide.docx",
        ),
    ]


def _remove_slide(prs, idx: int) -> None:
    slide_id = prs.slides._sldIdLst[idx]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[idx]


def _build_deck() -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.util import Inches as PptInches, Pt as PptPt

    colors = {
        "ink": PptRGB(0x29, 0x34, 0x36),
        "blue": PptRGB(0x0C, 0xA5, 0xF2),
        "teal": PptRGB(0x1F, 0xD9, 0xB8),
        "green": PptRGB(0x2A, 0xF5, 0x98),
        "gray1": PptRGB(0xF4, 0xF6, 0xF7),
        "gray2": PptRGB(0xD7, 0xDD, 0xE0),
        "gray3": PptRGB(0x73, 0x7A, 0x7B),
        "white": PptRGB(0xFF, 0xFF, 0xFF),
        "soft_blue": PptRGB(0xE8, 0xF7, 0xFE),
        "soft_teal": PptRGB(0xE9, 0xFB, 0xF7),
    }

    def remove_shape(shape) -> None:
        shape.element.getparent().remove(shape.element)

    def layout_by_name(prs, name: str, fallback_index: int = 0):
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
        return prs.slide_layouts[fallback_index]

    def blank_slide(prs, layout_name: str = "background-2"):
        slide = prs.slides.add_slide(layout_by_name(prs, layout_name))
        for shape in list(slide.shapes):
            if getattr(shape, "is_placeholder", False):
                remove_shape(shape)
        return slide

    def set_text(paragraph, text, size=14, color="ink", bold=False, align=None):
        paragraph.text = text
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Roboto"
            run.font.size = PptPt(size)
            run.font.bold = bold
            run.font.color.rgb = colors[color]

    def add_text(slide, text, left, top, width, height, size=14, color="ink", bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = PptPt(0)
        tf.margin_right = PptPt(0)
        set_text(tf.paragraphs[0], text, size=size, color=color, bold=bold, align=align)
        return box

    def add_header(slide, title, subtitle=None):
        logo = BRAND_IMAGES / "keycentrix-logo-from-template.png"
        if logo.exists():
            slide.shapes.add_picture(str(logo), PptInches(10.35), PptInches(0.22), width=PptInches(2.25))
        add_text(slide, title, 0.55, 0.28, 9.2, 0.42, size=22.5, bold=True)
        if subtitle:
            add_text(slide, subtitle, 0.56, 0.72, 9.4, 0.30, size=10.5, color="gray3")
        rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(0.56), PptInches(1.08), PptInches(1.35), PptInches(0.08))
        rule.fill.solid()
        rule.fill.fore_color.rgb = colors["blue"]
        rule.line.fill.background()

    def add_footer(slide):
        add_text(slide, "keycentrix | KCXDocumentor", 0.56, 7.05, 5.8, 0.2, size=8.5, color="gray3")

    def add_card(slide, left, top, width, height, heading, body, accent="blue", fill="gray1"):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[fill]
        shape.line.color.rgb = colors["white"]
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[accent]
        bar.line.fill.background()
        add_text(slide, heading, left + 0.2, top + 0.24, width - 0.4, 0.3, size=14, bold=True)
        box = slide.shapes.add_textbox(PptInches(left + 0.2), PptInches(top + 0.66), PptInches(width - 0.4), PptInches(height - 0.78))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = PptPt(0)
        tf.margin_right = PptPt(0)
        for idx, line in enumerate(body):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            set_text(p, line, size=10.5, color="ink")
            p.space_after = PptPt(5)

    def add_table(slide, left, top, width, height, headers, rows):
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = colors["white"]
        panel.line.color.rgb = colors["gray2"]
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(left), PptInches(top + 0.08), PptInches(width), PptInches(height - 0.08))
        table = table_shape.table
        for row_idx in range(len(rows) + 1):
            table.rows[row_idx].height = PptInches((height - 0.08) / (len(rows) + 1))
        for col_idx in range(len(headers)):
            table.columns[col_idx].width = PptInches(width / len(headers))
        for idx, header in enumerate(headers):
            cell = table.cell(0, idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["soft_blue"]
            cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            set_text(cell.text_frame.paragraphs[0], header, size=9.2, bold=True, align=PP_ALIGN.CENTER)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["gray1"] if col_idx == 0 else colors["white"]
                cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                set_text(cell.text_frame.paragraphs[0], value, size=8.8, bold=(col_idx == 0))

    prs = Presentation(str(PPT_TEMPLATE))
    for idx in range(len(prs.slides) - 1, 0, -1):
        _remove_slide(prs, idx)
    title_slide = prs.slides[0]
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "KCXDocumentor"
    add_text(title_slide, "Turning workflow recordings into reviewable user guides", 3.45, 4.82, 7.2, 0.42, size=15, color="gray3", align=PP_ALIGN.CENTER)
    add_text(title_slide, "Built with the DamienDev presentation recipe pattern", 3.45, 5.28, 7.2, 0.30, size=10.5, color="gray3", align=PP_ALIGN.CENTER)

    slides = [
        ("The Documentation Bottleneck", "BA, trainer, and implementation teams already capture useful walkthroughs.", [
            ("Current drag", ["Rewatch long recordings", "Identify workflow steps", "Capture screenshots", "Rewrite narration into user instructions"]),
            ("Why it matters", ["One-hour videos can take many hours to document", "Quality varies by author", "Manual screenshot selection slows publication"]),
        ]),
        ("What KCXDocumentor Changes", "The tool turns local recordings into reviewable DOCX guides.", [
            ("Guided workflow", ["Import recording and transcript", "Process locally", "Review screenshots", "Create and download DOCX"]),
            ("Quality gate", ["Rejected frames are excluded", "Reviewer notes travel into generation context", "Comments identify areas needing validation"]),
        ]),
        ("Local-First AI Boundary", "Raw video remains on the workstation.", [
            ("Local processing", ["FFmpeg, Whisper, OCR, and frame extraction run locally", "Artifacts stay in mapped workstation folders"]),
            ("AI payload", ["Compact transcript and metadata", "Reviewer-approved screenshots and notes", "No raw video sent to the provider"]),
        ]),
        ("Current Build Evidence", "The internal build is ready for hands-on validation.", [
            ("Application", ["Dockerized local web console", "Runtime Whisper setup", "Windows-friendly folder mapping"]),
            ("Cloud support", ["Entra sign-in with MSAL + PKCE", "Azure Function proxy", "Cosmos-backed AI Spend"]),
        ]),
        ("AI Spend Visibility", "Usage reporting makes the economics understandable.", []),
        ("Validation Path", "Run representative recordings through the same workflow the team will use.", [
            ("Inputs", ["Teams recordings with transcripts", "Recordings without transcripts using Whisper", "Multiple workflow lengths and styles"]),
            ("Measures", ["Time to first usable guide", "Screenshot relevance", "Reviewer edits", "Cost per guide and page"]),
        ]),
        ("Decision Request", "Approve internal use with BAs and trainers for focused validation.", [
            ("What approval enables", ["Representative workflow testing", "Quality feedback from target users", "Evidence for broader rollout planning"]),
            ("Success signal", ["Draft guides are readable", "Reviewers trust screenshot control", "Costs are visible and predictable"]),
        ]),
    ]

    for title, subtitle, cards in slides:
        slide = blank_slide(prs)
        add_header(slide, title, subtitle)
        add_footer(slide)
        if title == "AI Spend Visibility":
            add_table(
                slide,
                0.75,
                1.55,
                11.8,
                4.85,
                ["Metric", "Why Executives Care"],
                [
                    ["Documents", "Shows output volume."],
                    ["Tokens", "Explains model consumption."],
                    ["Estimated Cost", "Keeps spend transparent."],
                    ["Pages", "Connects cost to documentation produced."],
                    ["Cost Per Page", "Creates a comparable unit economics metric."],
                    ["User", "Attributes generation to the authenticated employee."],
                ],
            )
        else:
            left = 0.72
            top = 1.55
            width = 5.75
            for idx, (heading, body) in enumerate(cards):
                add_card(slide, left + idx * (width + 0.42), top, width, 4.7, heading, body, accent="teal" if idx else "blue", fill="soft_teal" if idx else "gray1")

    output = OUT_DIR / "KCXDocumentor - Executive Pitch Deck.pptx"
    prs.save(output)
    return output


def build_all() -> list[Path]:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    outputs = [_build_doc(spec) for spec in _doc_specs()]
    outputs.append(_build_deck())
    if SAMPLE_GUIDE.exists():
        sample_out = OUT_DIR / "KCXDocumentor - Sample Generated User Guide - Blink Rx.docx"
        shutil.copy2(SAMPLE_GUIDE, sample_out)
        _prepend_sample_recipe_cover(sample_out)
        outputs.append(sample_out)
    return outputs


def _prepend_sample_recipe_cover(path: Path) -> None:
    sample_doc = Document(path)
    sample_doc.core_properties.title = "KCXDocumentor - Sample Generated User Guide - Blink Rx"
    sample_doc.core_properties.author = "keycentrix"
    sample_doc.core_properties.comments = "Showcase copy packaged with the DamienDev document recipe pattern."

    front = Document()
    spec = DocSpec(
        title="KCXDocumentor - Sample Generated User Guide - Blink Rx",
        subtitle="Example generated guide packaged for executive review",
        document_id="KCXDOC-SAMPLE-001",
        status="Internal showcase sample",
        owner="keycentrix Product and Training",
        audience="Executive Team, Product, Training, Implementation, and Documentation Reviewers",
        purpose="Show a real KCXDocumentor generated guide as a packaged sample artifact.",
        scope="Includes recipe front matter followed by the generated Blink Rx guide output.",
        sections=[],
        output_name=path.name,
    )
    _style_document(front)
    _header_footer(front, spec.title)
    _cover(front, spec)
    front.add_heading("Package Notes", level=1)
    _bullets(
        front,
        [
            "The following pages are the generated guide output used to demonstrate KCXDocumentor capability.",
            "This front matter was added for the executive collateral package.",
            "The generated guide content remains available for review behind this package section.",
        ],
    )
    front.add_page_break()

    front_body = front.element.body
    sample_body = sample_doc.element.body
    for element in reversed(list(front_body)):
        if element.tag.endswith("sectPr"):
            continue
        sample_body.insert(0, deepcopy(element))
    sample_doc.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build KCXDocumentor prototype collateral using the DamienDev recipe pattern.")
    parser.add_argument("--all", action="store_true", help="Build the full prototype collateral package.")
    args = parser.parse_args()
    outputs = build_all() if args.all else build_all()
    for output in outputs:
        print(f"Built {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
